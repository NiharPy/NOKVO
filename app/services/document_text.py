"""Plain-text extraction from uploaded documents (PDF / DOCX / PPTX / TXT).

Shared by the outbound-campaign reference-doc upload and the real-estate
Brochure Analyzer. Best-effort: an unsupported / unparseable file degrades to a
utf-8 decode rather than raising, so a bad upload never 500s the caller.
"""
from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)


def extract_document_text(filename: str, content: bytes) -> str:
    """Extract plain text from a PDF, DOCX, or TXT file (by extension).

    Binary formats (PDF/DOCX) return "" when their parser is missing OR when the
    file carries no extractable text (e.g. a scanned / image-only PDF) — we
    deliberately do NOT utf-8-decode binary bytes, since that yields garbage that
    silently poisons downstream LLM extraction. Plain-text inputs decode."""
    ext = (filename or "").rsplit(".", 1)[-1].lower()
    if ext == "txt":
        return content.decode("utf-8", errors="replace")
    if ext == "pdf":
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            if text.strip():
                return text
        except ImportError:
            pass
        except Exception:
            logger.warning("document_text: pypdf failed to read PDF", exc_info=True)
        try:
            import pdfminer.high_level as pm

            return pm.extract_text(io.BytesIO(content)) or ""
        except ImportError:
            logger.warning("document_text: no PDF parser installed (pypdf/pdfminer.six)")
            return ""
        except Exception:
            logger.warning("document_text: pdfminer failed to read PDF", exc_info=True)
            return ""
    if ext in ("docx", "doc"):
        try:
            import docx

            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            logger.warning("document_text: python-docx not installed")
            return ""
        except Exception:
            logger.warning("document_text: failed to read DOCX", exc_info=True)
            return ""
    if ext == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                parts.append(line)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                parts.append(" | ".join(cells))
            return "\n".join(parts)
        except ImportError:
            logger.warning("document_text: python-pptx not installed")
            return ""
        except Exception:
            logger.warning("document_text: failed to read PPTX", exc_info=True)
            return ""
    return content.decode("utf-8", errors="replace")
