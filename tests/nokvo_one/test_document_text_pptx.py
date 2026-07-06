"""PPTX text extraction — generated fixture deck through extract_document_text."""
import io

from pptx import Presentation
from pptx.util import Inches

from app.services.document_text import extract_document_text


def _make_deck() -> bytes:
    prs = Presentation()
    # Title slide.
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Skyline Heights — Premium 3BHK Flats"
    slide.placeholders[1].text = "Kokapet, Hyderabad · From ₹1.2 Cr"
    # Content slide with a bullet list.
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Why buyers choose us"
    body = slide2.placeholders[1].text_frame
    body.text = "RERA approved"
    body.add_paragraph().text = "Possession Dec 2027"
    # Slide with a table.
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide3.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "Config"
    table.cell(0, 1).text = "Price"
    table.cell(1, 0).text = "3BHK"
    table.cell(1, 1).text = "1.2 Cr"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_extracts_titles_bullets_and_tables():
    text = extract_document_text("pitch.pptx", _make_deck())
    assert "Skyline Heights" in text
    assert "Kokapet" in text
    assert "RERA approved" in text
    assert "Possession Dec 2027" in text
    assert "3BHK | 1.2 Cr" in text  # table rows join with |


def test_pptx_garbage_degrades_to_empty():
    assert extract_document_text("bad.pptx", b"not a pptx at all") == ""


def test_other_formats_untouched():
    assert extract_document_text("note.txt", "hello".encode()) == "hello"
